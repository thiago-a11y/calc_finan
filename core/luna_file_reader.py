"""
Luna File Reader — Extrator de conteúdo de arquivos anexados.

Lê o conteúdo real dos arquivos enviados pelo usuário para que
a Luna possa analisar e responder com base nos dados reais.

Formatos suportados:
- XLSX/XLS: Extrai todas as abas como texto tabular
- DOCX: Extrai texto completo (parágrafos + tabelas)
- PDF: Extrai texto de todas as páginas
- CSV/TSV: Lê como texto tabular
- TXT/MD/JSON/HTML/XML: Lê diretamente
- Imagens (JPG/PNG): Descreve como imagem (não OCR por enquanto)
- PPTX: Extrai texto dos slides

Limite: trunca conteúdo em 15.000 caracteres para não estourar contexto.
"""

import os
import csv
import json
import logging
from io import BytesIO, StringIO
from pathlib import Path

logger = logging.getLogger("luna.file_reader")

# Limite de caracteres por arquivo no contexto do LLM
LIMITE_CARACTERES = 15_000

# Diretório base de uploads
UPLOAD_DIR = Path(__file__).parent.parent / "uploads" / "chat"


def extrair_conteudo_arquivo(url: str, nome_original: str, tipo: str) -> str | None:
    """
    Extrai o conteúdo textual de um arquivo anexado.

    Args:
        url: URL relativa do arquivo (ex: /uploads/chat/luna_abc123.xlsx)
        nome_original: Nome original do arquivo enviado pelo usuário
        tipo: MIME type do arquivo

    Returns:
        String com o conteúdo extraído, ou None se não conseguir ler.
    """
    try:
        # Resolver caminho do arquivo no disco
        caminho = _resolver_caminho(url)
        if not caminho or not caminho.exists():
            logger.warning(f"Arquivo não encontrado: {url}")
            return None

        extensao = Path(nome_original).suffix.lower()

        # Escolher extrator baseado na extensão
        extratores = {
            ".xlsx": _ler_xlsx,
            ".xls": _ler_xlsx,
            ".csv": _ler_csv,
            ".tsv": _ler_tsv,
            ".docx": _ler_docx,
            ".doc": _ler_docx,
            ".pdf": _ler_pdf,
            ".pptx": _ler_pptx,
            ".txt": _ler_texto,
            ".md": _ler_texto,
            ".json": _ler_json,
            ".html": _ler_texto,
            ".htm": _ler_texto,
            ".xml": _ler_texto,
            ".py": _ler_texto,
            ".js": _ler_texto,
            ".ts": _ler_texto,
            ".jsx": _ler_texto,
            ".tsx": _ler_texto,
            ".css": _ler_texto,
            ".sql": _ler_texto,
            ".yaml": _ler_texto,
            ".yml": _ler_texto,
            ".env": _ler_texto,
            ".log": _ler_texto,
        }

        extrator = extratores.get(extensao)
        if not extrator:
            # Imagens e outros binários
            if tipo and tipo.startswith("image/"):
                return f"[Imagem: {nome_original} — análise visual não disponível]"
            return f"[Arquivo binário: {nome_original} — formato não suportado para leitura]"

        conteudo = extrator(caminho)

        if conteudo and len(conteudo) > LIMITE_CARACTERES:
            conteudo = conteudo[:LIMITE_CARACTERES] + f"\n\n... [conteúdo truncado — total: {len(conteudo):,} caracteres]"

        return conteudo

    except Exception as e:
        logger.error(f"Erro ao extrair conteúdo de {nome_original}: {e}")
        return f"[Erro ao ler {nome_original}: {str(e)}]"


def _resolver_caminho(url: str) -> Path | None:
    """Resolve a URL relativa para um caminho absoluto no disco."""
    if not url:
        return None

    # URL tipo /uploads/chat/luna_abc123.xlsx
    if url.startswith("/uploads/chat/"):
        nome_arquivo = url.split("/")[-1]
        return UPLOAD_DIR / nome_arquivo

    # Se for caminho absoluto
    if url.startswith("/") and os.path.exists(url):
        return Path(url)

    return None


def _ler_xlsx(caminho: Path) -> str:
    """Extrai conteúdo de planilha Excel."""
    try:
        import openpyxl
    except ImportError:
        return "[Erro: biblioteca openpyxl não instalada]"

    wb = openpyxl.load_workbook(str(caminho), read_only=True, data_only=True)
    partes = []

    for nome_aba in wb.sheetnames:
        ws = wb[nome_aba]
        partes.append(f"\n📊 Aba: {nome_aba}")
        partes.append("-" * 50)

        linhas_processadas = 0
        for row in ws.iter_rows(values_only=True):
            valores = []
            for cell in row:
                if cell is None:
                    valores.append("")
                else:
                    valores.append(str(cell))
            partes.append(" | ".join(valores))
            linhas_processadas += 1

            # Limitar linhas por aba
            if linhas_processadas >= 500:
                partes.append(f"... [mais linhas omitidas — total estimado na aba]")
                break

        partes.append("")

    wb.close()
    return "\n".join(partes)


def _ler_csv(caminho: Path) -> str:
    """Extrai conteúdo de CSV."""
    partes = ["📊 Dados CSV:"]
    partes.append("-" * 50)

    with open(caminho, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            partes.append(" | ".join(row))
            if i >= 500:
                partes.append("... [mais linhas omitidas]")
                break

    return "\n".join(partes)


def _ler_tsv(caminho: Path) -> str:
    """Extrai conteúdo de TSV."""
    partes = ["📊 Dados TSV:"]
    partes.append("-" * 50)

    with open(caminho, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter="\t")
        for i, row in enumerate(reader):
            partes.append(" | ".join(row))
            if i >= 500:
                partes.append("... [mais linhas omitidas]")
                break

    return "\n".join(partes)


def _ler_docx(caminho: Path) -> str:
    """Extrai texto de documento Word."""
    try:
        from docx import Document
    except ImportError:
        return "[Erro: biblioteca python-docx não instalada]"

    doc = Document(str(caminho))
    partes = ["📄 Documento Word:"]
    partes.append("-" * 50)

    # Parágrafos
    for para in doc.paragraphs:
        if para.text.strip():
            partes.append(para.text)

    # Tabelas
    for i, table in enumerate(doc.tables):
        partes.append(f"\n📊 Tabela {i + 1}:")
        for row in table.rows:
            valores = [cell.text.strip() for cell in row.cells]
            partes.append(" | ".join(valores))

    return "\n".join(partes)


def _ler_pdf(caminho: Path) -> str:
    """Extrai texto de PDF."""
    try:
        from reportlab.lib.pagesizes import letter  # noqa: F401
    except ImportError:
        pass

    # Tentar com PyPDF2 primeiro
    try:
        import pypdf
        reader = pypdf.PdfReader(str(caminho))
        partes = [f"📄 PDF ({len(reader.pages)} páginas):"]
        partes.append("-" * 50)

        for i, page in enumerate(reader.pages):
            texto = page.extract_text()
            if texto and texto.strip():
                partes.append(f"\n--- Página {i + 1} ---")
                partes.append(texto.strip())

        return "\n".join(partes)
    except ImportError:
        pass

    # Tentar com PyMuPDF
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(caminho))
        partes = [f"📄 PDF ({len(doc)} páginas):"]
        partes.append("-" * 50)

        for i, page in enumerate(doc):
            texto = page.get_text()
            if texto and texto.strip():
                partes.append(f"\n--- Página {i + 1} ---")
                partes.append(texto.strip())

        doc.close()
        return "\n".join(partes)
    except ImportError:
        pass

    return "[Erro: nenhuma biblioteca PDF disponível — instale pypdf ou PyMuPDF]"


def _ler_pptx(caminho: Path) -> str:
    """Extrai texto de apresentação PowerPoint."""
    try:
        from pptx import Presentation
    except ImportError:
        return "[Erro: biblioteca python-pptx não instalada]"

    prs = Presentation(str(caminho))
    partes = [f"📊 Apresentação ({len(prs.slides)} slides):"]
    partes.append("-" * 50)

    for i, slide in enumerate(prs.slides):
        partes.append(f"\n--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                partes.append(shape.text)

    return "\n".join(partes)


def _ler_texto(caminho: Path) -> str:
    """Lê arquivo de texto puro."""
    with open(caminho, "r", encoding="utf-8", errors="replace") as f:
        conteudo = f.read()

    extensao = caminho.suffix.lower()
    icone = "📄" if extensao in (".txt", ".md") else "💻"
    return f"{icone} Conteúdo de {caminho.name}:\n{'-' * 50}\n{conteudo}"


def _ler_json(caminho: Path) -> str:
    """Lê e formata JSON."""
    with open(caminho, "r", encoding="utf-8", errors="replace") as f:
        dados = json.load(f)

    formatado = json.dumps(dados, indent=2, ensure_ascii=False)
    return f"📋 JSON:\n{'-' * 50}\n{formatado}"
