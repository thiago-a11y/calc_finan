"""
Luna File Generator — Gerador de arquivos para a assistente Luna.

Gera arquivos nos formatos: XLSX, DOCX, PPTX, PDF, TXT, MD, CSV, JSON, HTML
Os arquivos são salvos no diretório de uploads e servidos via URL.

Uso:
    from core.luna_file_generator import gerar_arquivo
    resultado = gerar_arquivo("xlsx", conteudo, "planilha_vendas", usuario_nome)
    # {"url": "/uploads/chat/abc123.xlsx", "nome": "planilha_vendas.xlsx", ...}
"""

import os
import json
import uuid
import logging
from datetime import datetime
from io import BytesIO

logger = logging.getLogger("synerium.luna.files")

UPLOAD_DIR = os.path.expanduser("~/synerium-factory/data/uploads/chat")
# No servidor AWS
if os.path.exists("/opt/synerium-factory"):
    UPLOAD_DIR = "/opt/synerium-factory/data/uploads/chat"


def _gerar_nome_unico(extensao: str) -> str:
    """Gera nome único para o arquivo."""
    return f"luna_{uuid.uuid4().hex[:10]}{extensao}"


def gerar_arquivo(
    formato: str,
    conteudo: str,
    nome_base: str = "arquivo",
    usuario_nome: str = "",
    titulo: str = "",
) -> dict:
    """
    Gera um arquivo no formato especificado.

    Args:
        formato: Tipo do arquivo (xlsx, docx, pptx, pdf, txt, md, csv, json, html)
        conteudo: Conteúdo a ser convertido no arquivo
        nome_base: Nome base do arquivo (sem extensão)
        usuario_nome: Nome do usuário que solicitou
        titulo: Título do documento (opcional)

    Returns:
        Dict com url, nome, tipo, tamanho, formato
    """
    os.makedirs(UPLOAD_DIR, exist_ok=True)

    formato = formato.lower().strip().lstrip(".")

    geradores = {
        "xlsx": _gerar_xlsx,
        "docx": _gerar_docx,
        "pptx": _gerar_pptx,
        "pdf": _gerar_pdf,
        "txt": _gerar_txt,
        "md": _gerar_md,
        "csv": _gerar_csv,
        "json": _gerar_json,
        "html": _gerar_html,
    }

    gerador = geradores.get(formato)
    if not gerador:
        raise ValueError(f"Formato '{formato}' não suportado. Use: {', '.join(geradores.keys())}")

    extensao = f".{formato}"
    nome_arquivo = _gerar_nome_unico(extensao)
    caminho = os.path.join(UPLOAD_DIR, nome_arquivo)

    try:
        gerador(conteudo, caminho, titulo or nome_base)
        tamanho = os.path.getsize(caminho)

        logger.info(f"[Luna Files] Gerado: {nome_base}{extensao} ({tamanho} bytes) por {usuario_nome}")

        return {
            "url": f"/uploads/chat/{nome_arquivo}",
            "nome": f"{nome_base}{extensao}",
            "nome_salvo": nome_arquivo,
            "formato": formato,
            "tipo": _tipo_por_formato(formato),
            "tamanho": tamanho,
            "gerado_por": "Luna",
            "gerado_em": datetime.now().isoformat(),
        }
    except Exception as e:
        logger.error(f"[Luna Files] Erro ao gerar {formato}: {e}")
        raise


def _tipo_por_formato(formato: str) -> str:
    """Retorna o tipo de arquivo pelo formato."""
    tipos = {
        "xlsx": "planilha",
        "csv": "planilha",
        "docx": "documento",
        "pptx": "apresentacao",
        "pdf": "pdf",
        "txt": "texto",
        "md": "texto",
        "json": "dados",
        "html": "html",
    }
    return tipos.get(formato, "documento")


# =====================================================================
# Geradores por formato
# =====================================================================


def _gerar_xlsx(conteudo: str, caminho: str, titulo: str):
    """Gera planilha XLSX a partir de conteúdo texto/tabular."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = titulo[:31]  # Max 31 chars para nome da aba

    # Estilos
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="10B981", end_color="10B981", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    # Parsear conteúdo (tentativas: tabela markdown, CSV, linhas)
    linhas = conteudo.strip().split("\n")

    for row_idx, linha in enumerate(linhas, 1):
        # Pular separadores de tabela markdown (|---|---|)
        if linha.strip().startswith("|") and set(linha.replace("|", "").replace("-", "").replace(":", "").strip()) == set():
            continue

        # Separar por | (tabela markdown) ou \t (TSV) ou , (CSV)
        if "|" in linha:
            celulas = [c.strip() for c in linha.strip().strip("|").split("|")]
        elif "\t" in linha:
            celulas = linha.split("\t")
        elif "," in linha and row_idx <= 1:
            celulas = [c.strip().strip('"') for c in linha.split(",")]
        else:
            celulas = [linha]

        for col_idx, valor in enumerate(celulas, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=valor)
            cell.border = thin_border
            cell.alignment = Alignment(vertical="center")

            # Tentar converter números
            try:
                cell.value = float(valor.replace(",", "."))
                cell.number_format = "#,##0.00"
            except (ValueError, AttributeError):
                pass

            # Estilo do header (primeira linha)
            if row_idx == 1:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_align

    # Auto-ajustar largura das colunas
    for col in ws.columns:
        max_length = 0
        col_letter = col[0].column_letter
        for cell in col:
            if cell.value:
                max_length = max(max_length, len(str(cell.value)))
        ws.column_dimensions[col_letter].width = min(max_length + 4, 50)

    wb.save(caminho)


def _gerar_docx(conteudo: str, caminho: str, titulo: str):
    """Gera documento Word DOCX a partir de conteúdo markdown/texto."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Estilo do título
    title = doc.add_heading(titulo, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadados
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(f"Gerado por Luna — Synerium Factory | {datetime.now().strftime('%d/%m/%Y %H:%M')}")
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(128, 128, 128)

    doc.add_paragraph()  # Espaço

    # Parsear conteúdo markdown básico
    linhas = conteudo.strip().split("\n")
    i = 0
    while i < len(linhas):
        linha = linhas[i]

        # Headers
        if linha.startswith("### "):
            doc.add_heading(linha[4:], level=3)
        elif linha.startswith("## "):
            doc.add_heading(linha[3:], level=2)
        elif linha.startswith("# "):
            doc.add_heading(linha[2:], level=1)
        # Lista
        elif linha.strip().startswith("- ") or linha.strip().startswith("* "):
            texto = linha.strip().lstrip("-* ").strip()
            doc.add_paragraph(texto, style="List Bullet")
        elif linha.strip() and linha.strip()[0].isdigit() and ". " in linha:
            texto = linha.strip().split(". ", 1)[1] if ". " in linha else linha
            doc.add_paragraph(texto, style="List Number")
        # Bloco de código
        elif linha.strip().startswith("```"):
            codigo = []
            i += 1
            while i < len(linhas) and not linhas[i].strip().startswith("```"):
                codigo.append(linhas[i])
                i += 1
            if codigo:
                p = doc.add_paragraph()
                run = p.add_run("\n".join(codigo))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
        # Parágrafo normal
        elif linha.strip():
            p = doc.add_paragraph()
            # Processar negrito (**texto**)
            partes = linha.split("**")
            for j, parte in enumerate(partes):
                run = p.add_run(parte)
                if j % 2 == 1:  # Índice ímpar = dentro de **
                    run.bold = True
        else:
            doc.add_paragraph()  # Linha vazia

        i += 1

    doc.save(caminho)


def _gerar_pptx(conteudo: str, caminho: str, titulo: str):
    """Gera apresentação PowerPoint PPTX."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()

    # Slide de título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = f"Gerado por Luna — {datetime.now().strftime('%d/%m/%Y')}"

    # Parsear slides (separados por ## ou ---)
    linhas = conteudo.strip().split("\n")
    slide_atual = []
    titulo_slide = ""

    for linha in linhas:
        if linha.startswith("## ") or linha.strip() == "---":
            # Salvar slide anterior
            if titulo_slide or slide_atual:
                _adicionar_slide_pptx(prs, titulo_slide, slide_atual)
                slide_atual = []

            if linha.startswith("## "):
                titulo_slide = linha[3:].strip()
            else:
                titulo_slide = ""
        else:
            if linha.strip():
                slide_atual.append(linha)

    # Último slide
    if titulo_slide or slide_atual:
        _adicionar_slide_pptx(prs, titulo_slide, slide_atual)

    prs.save(caminho)


def _adicionar_slide_pptx(prs, titulo: str, linhas: list):
    """Adiciona um slide ao PPTX."""
    from pptx.util import Pt

    slide = prs.slides.add_slide(prs.slide_layouts[1])

    if titulo:
        slide.shapes.title.text = titulo

    if linhas and slide.placeholders.get(1):
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, linha in enumerate(linhas):
            texto = linha.strip().lstrip("-*• ").strip()
            if i == 0:
                tf.paragraphs[0].text = texto
            else:
                p = tf.add_paragraph()
                p.text = texto
            tf.paragraphs[-1].font.size = Pt(18)


def _gerar_pdf(conteudo: str, caminho: str, titulo: str):
    """Gera PDF a partir de conteúdo texto."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    doc = SimpleDocTemplate(caminho, pagesize=A4,
                           topMargin=2*cm, bottomMargin=2*cm,
                           leftMargin=2.5*cm, rightMargin=2.5*cm)

    styles = getSampleStyleSheet()

    # Estilo customizado
    styles.add(ParagraphStyle(
        "LunaTitle",
        parent=styles["Title"],
        textColor=HexColor("#10B981"),
        fontSize=20,
        spaceAfter=12,
    ))
    styles.add(ParagraphStyle(
        "LunaMeta",
        parent=styles["Normal"],
        textColor=HexColor("#888888"),
        fontSize=9,
        alignment=1,  # Center
        spaceAfter=20,
    ))

    elementos = []

    # Título
    elementos.append(Paragraph(titulo, styles["LunaTitle"]))
    elementos.append(Paragraph(
        f"Gerado por Luna — Synerium Factory | {datetime.now().strftime('%d/%m/%Y %H:%M')}",
        styles["LunaMeta"]
    ))
    elementos.append(Spacer(1, 12))

    # Conteúdo
    for linha in conteudo.strip().split("\n"):
        if linha.startswith("### "):
            elementos.append(Paragraph(linha[4:], styles["Heading3"]))
        elif linha.startswith("## "):
            elementos.append(Paragraph(linha[3:], styles["Heading2"]))
        elif linha.startswith("# "):
            elementos.append(Paragraph(linha[2:], styles["Heading1"]))
        elif linha.strip().startswith("- ") or linha.strip().startswith("* "):
            texto = linha.strip().lstrip("-* ").strip()
            elementos.append(Paragraph(f"• {texto}", styles["Normal"]))
        elif linha.strip():
            # Converter **negrito** para <b>
            texto = linha.replace("**", "<b>", 1)
            while "**" in texto:
                texto = texto.replace("**", "</b>", 1)
                if "**" in texto:
                    texto = texto.replace("**", "<b>", 1)
            elementos.append(Paragraph(texto, styles["Normal"]))
        else:
            elementos.append(Spacer(1, 6))

    doc.build(elementos)


def _gerar_txt(conteudo: str, caminho: str, titulo: str):
    """Gera arquivo TXT simples."""
    with open(caminho, "w", encoding="utf-8") as f:
        f.write(f"{titulo}\n{'='*len(titulo)}\n\n")
        f.write(conteudo)
        f.write(f"\n\n---\nGerado por Luna — Synerium Factory | {datetime.now().strftime('%d/%m/%Y %H:%M')}")


def _gerar_md(conteudo: str, caminho: str, titulo: str):
    """Gera arquivo Markdown."""
    with open(caminho, "w", encoding="utf-8") as f:
        f.write(f"# {titulo}\n\n")
        f.write(conteudo)
        f.write(f"\n\n---\n*Gerado por Luna — Synerium Factory | {datetime.now().strftime('%d/%m/%Y %H:%M')}*")


def _gerar_csv(conteudo: str, caminho: str, titulo: str):
    """Gera arquivo CSV."""
    import csv

    linhas = conteudo.strip().split("\n")

    with open(caminho, "w", encoding="utf-8-sig", newline="") as f:
        writer = csv.writer(f)
        for linha in linhas:
            # Parsear tabela markdown ou CSV
            if "|" in linha:
                celulas = [c.strip() for c in linha.strip().strip("|").split("|")]
                # Pular separadores
                if all(set(c.replace("-", "").replace(":", "").strip()) == set() for c in celulas if c):
                    continue
                writer.writerow(celulas)
            elif "," in linha:
                writer.writerow([c.strip().strip('"') for c in linha.split(",")])
            elif "\t" in linha:
                writer.writerow(linha.split("\t"))
            else:
                writer.writerow([linha])


def _gerar_json(conteudo: str, caminho: str, titulo: str):
    """Gera arquivo JSON."""
    with open(caminho, "w", encoding="utf-8") as f:
        # Tentar parsear como JSON válido
        try:
            data = json.loads(conteudo)
            json.dump(data, f, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            # Se não for JSON, criar estrutura
            json.dump({
                "titulo": titulo,
                "conteudo": conteudo,
                "gerado_por": "Luna",
                "gerado_em": datetime.now().isoformat(),
            }, f, ensure_ascii=False, indent=2)


def _gerar_html(conteudo: str, caminho: str, titulo: str):
    """Gera arquivo HTML."""
    html = f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{titulo}</title>
    <style>
        body {{ font-family: 'Inter', -apple-system, sans-serif; max-width: 800px; margin: 0 auto; padding: 2rem; color: #1a1a1a; line-height: 1.6; }}
        h1 {{ color: #10B981; border-bottom: 2px solid #10B981; padding-bottom: 0.5rem; }}
        h2 {{ color: #333; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1rem 0; }}
        th {{ background: #10B981; color: white; padding: 0.75rem; text-align: left; }}
        td {{ border: 1px solid #ddd; padding: 0.75rem; }}
        tr:nth-child(even) {{ background: #f9f9f9; }}
        code {{ background: #f4f4f4; padding: 0.2rem 0.4rem; border-radius: 3px; font-size: 0.9em; }}
        pre {{ background: #1a1b26; color: #e0e0e0; padding: 1rem; border-radius: 8px; overflow-x: auto; }}
        .footer {{ margin-top: 2rem; padding-top: 1rem; border-top: 1px solid #eee; color: #888; font-size: 0.8rem; }}
    </style>
</head>
<body>
    <h1>{titulo}</h1>
    {conteudo}
    <div class="footer">Gerado por Luna — Synerium Factory | {datetime.now().strftime('%d/%m/%Y %H:%M')}</div>
</body>
</html>"""

    with open(caminho, "w", encoding="utf-8") as f:
        f.write(html)
